#!/usr/bin/env python3
import argparse
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


SVG_NS = "{http://www.w3.org/2000/svg}"


def strip_ns(tag):
    return tag.split("}", 1)[-1]


def num(value, default=0.0):
    if value is None:
        return default
    m = re.match(r"[-+]?\d*\.?\d+", str(value))
    return float(m.group(0)) if m else default


def color(value, default="000000"):
    if not value or value == "none":
        return None
    value = value.strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    if len(value) != 6:
        value = default
    return RGBColor.from_string(value.upper())


def parse_points(value):
    nums = [float(x) for x in re.findall(r"[-+]?\d*\.?\d+", value or "")]
    return list(zip(nums[0::2], nums[1::2]))


def bbox_from_points(points):
    xs = [p[0] for p in points]
    ys = [p[1] for p in points]
    return min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)


def is_star(points):
    return len(points) in (10, 11)


def set_line_style(line, elem):
    stroke = color(elem.attrib.get("stroke"), "000000")
    if stroke:
        line.line.color.rgb = stroke
    line.line.width = Pt(num(elem.attrib.get("stroke-width"), 1))


def main():
    ap = argparse.ArgumentParser(description="Convert a simple SVG component to editable PPTX shapes.")
    ap.add_argument("svg")
    ap.add_argument("pptx")
    ap.add_argument("--slide-width", type=float, default=13.333333)
    ap.add_argument("--slide-height", type=float, default=7.5)
    args = ap.parse_args()

    root = ET.parse(args.svg).getroot()
    view_box = root.attrib.get("viewBox")
    if view_box:
        _, _, vb_w, vb_h = [float(x) for x in view_box.replace(",", " ").split()]
    else:
        vb_w = num(root.attrib.get("width"), 1920)
        vb_h = num(root.attrib.get("height"), 1080)

    prs = Presentation()
    prs.slide_width = Inches(args.slide_width)
    prs.slide_height = Inches(args.slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def emu_x(x):
        return Inches(x / vb_w * args.slide_width)

    def emu_y(y):
        return Inches(y / vb_h * args.slide_height)

    def set_style(shape, elem):
        fill = color(elem.attrib.get("fill"), None)
        stroke = color(elem.attrib.get("stroke"), "000000")
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        if stroke is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = stroke
            sw = num(elem.attrib.get("stroke-width"), 1)
            shape.line.width = Pt(sw)

    for elem in root.iter():
        tag = strip_ns(elem.tag)
        if tag == "svg":
            continue
        if tag == "rect":
            x = num(elem.attrib.get("x"))
            y = num(elem.attrib.get("y"))
            w = num(elem.attrib.get("width"))
            h = num(elem.attrib.get("height"))
            rx = num(elem.attrib.get("rx"))
            shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rx else MSO_AUTO_SHAPE_TYPE.RECTANGLE
            shp = slide.shapes.add_shape(shape_type, emu_x(x), emu_y(y), emu_x(w), emu_y(h))
            set_style(shp, elem)
        elif tag == "circle":
            cx = num(elem.attrib.get("cx"))
            cy = num(elem.attrib.get("cy"))
            r = num(elem.attrib.get("r"))
            shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu_x(cx - r), emu_y(cy - r), emu_x(2 * r), emu_y(2 * r))
            set_style(shp, elem)
        elif tag == "ellipse":
            cx = num(elem.attrib.get("cx"))
            cy = num(elem.attrib.get("cy"))
            rx = num(elem.attrib.get("rx"))
            ry = num(elem.attrib.get("ry"))
            shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu_x(cx - rx), emu_y(cy - ry), emu_x(2 * rx), emu_y(2 * ry))
            set_style(shp, elem)
        elif tag == "line":
            x1 = num(elem.attrib.get("x1"))
            y1 = num(elem.attrib.get("y1"))
            x2 = num(elem.attrib.get("x2"))
            y2 = num(elem.attrib.get("y2"))
            line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, emu_x(x1), emu_y(y1), emu_x(x2), emu_y(y2))
            set_line_style(line, elem)
        elif tag in ("polygon", "polyline"):
            pts = parse_points(elem.attrib.get("points"))
            if not pts:
                continue
            if tag == "polyline":
                for p1, p2 in zip(pts, pts[1:]):
                    line = slide.shapes.add_connector(
                        MSO_CONNECTOR_TYPE.STRAIGHT,
                        emu_x(p1[0]),
                        emu_y(p1[1]),
                        emu_x(p2[0]),
                        emu_y(p2[1]),
                    )
                    set_line_style(line, elem)
            elif tag == "polygon" and is_star(pts):
                x, y, w, h = bbox_from_points(pts)
                shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, emu_x(x), emu_y(y), emu_x(w), emu_y(h))
                set_style(shp, elem)
            else:
                raise ValueError("unsupported polygon; simplify SVG, use ppt-master full converter, or reclassify as PNG")

    prs.save(args.pptx)
    print(args.pptx)


if __name__ == "__main__":
    main()
