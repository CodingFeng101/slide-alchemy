#!/usr/bin/env python3
import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SHAPES = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
    "star": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
}

ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def rgb(value, default="000000"):
    if value is None:
        value = default
    value = str(value).strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    return RGBColor.from_string(value.upper())


def as_bool(value, default=False):
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in ("1", "true", "yes", "y")
    return bool(value)


def resolve(base_dir, path):
    p = Path(path)
    return p if p.is_absolute() else base_dir / p


def bbox(item, ref_w, ref_h, slide_w, slide_h):
    if "bbox_frac" in item:
        x, y, w, h = item["bbox_frac"]
    else:
        x, y, w, h = item["bbox_px"]
        x, y, w, h = x / ref_w, y / ref_h, w / ref_w, h / ref_h
    return Inches(x * slide_w), Inches(y * slide_h), Inches(w * slide_w), Inches(h * slide_h)


def add_geometry(slide, item, ref_w, ref_h, slide_w, slide_h):
    kind = item["kind"]
    if kind == "line":
        x, y, w, h = bbox(item, ref_w, ref_h, slide_w, slide_h)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, max(h, Pt(1)))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(item.get("fill") or item.get("line"), "C00000")
        shape.line.fill.background()
        return shape

    if kind not in SHAPES:
        raise ValueError(f"unsupported geometry kind: {kind}")

    shape = slide.shapes.add_shape(SHAPES[kind], *bbox(item, ref_w, ref_h, slide_w, slide_h))
    fill = item.get("fill")
    if fill is None or fill == "none":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if "fill_transparency" in item:
            shape.fill.transparency = int(item["fill_transparency"])

    line = item.get("line")
    if line is None or line == "none":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(float(item.get("line_width", 1)))
    return shape


def add_text(slide, item, ref_w, ref_h, slide_w, slide_h):
    box = slide.shapes.add_textbox(*bbox(item, ref_w, ref_h, slide_w, slide_h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = ALIGN.get(item.get("align", "center"), PP_ALIGN.CENTER)
    run = p.add_run()
    run.text = item.get("text", "")
    run.font.name = item.get("font_family", "Microsoft YaHei")
    run.font.size = Pt(float(item.get("font_size_pt", 18)))
    run.font.bold = as_bool(item.get("bold"), False)
    run.font.italic = as_bool(item.get("italic"), False)
    run.font.color.rgb = rgb(item.get("color"), "000000")
    return box


def compose(spec_path, out_path):
    spec_path = Path(spec_path)
    base_dir = spec_path.parent
    spec = json.loads(spec_path.read_text(encoding="utf-8-sig"))
    ref_w = spec.get("ref_width", 1920)
    ref_h = spec.get("ref_height", 1080)
    slide_w = float(spec.get("slide_width_in", 13.333333))
    slide_h = float(spec.get("slide_height_in", 7.5))

    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    blank = prs.slide_layouts[6]

    for slide_spec in spec["slides"]:
        slide = prs.slides.add_slide(blank)
        if slide_spec.get("base"):
            slide.shapes.add_picture(str(resolve(base_dir, slide_spec["base"])), 0, 0, prs.slide_width, prs.slide_height)

        for item in slide_spec.get("geometry", []):
            add_geometry(slide, item, ref_w, ref_h, slide_w, slide_h)

        for item in slide_spec.get("images", []):
            slide.shapes.add_picture(str(resolve(base_dir, item["path"])), *bbox(item, ref_w, ref_h, slide_w, slide_h))

        for item in slide_spec.get("texts", []):
            add_text(slide, item, ref_w, ref_h, slide_w, slide_h)

    prs.save(out_path)
    return out_path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("compose_json")
    ap.add_argument("output_pptx")
    args = ap.parse_args()
    print(compose(args.compose_json, args.output_pptx))


if __name__ == "__main__":
    main()
